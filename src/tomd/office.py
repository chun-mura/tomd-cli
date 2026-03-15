"""Office document (docx, pptx, xlsx) post-processing for MarkItDown output."""

from __future__ import annotations

import re
from pathlib import Path


def _correct_docx_headings(src: Path, text: str) -> str:
    """Correct heading levels in MarkItDown output using python-docx styles.

    Reads the .docx paragraph styles to build a map of heading text -> level,
    then ensures the Markdown output uses the correct heading markers.
    """
    try:
        from docx import Document  # type: ignore[import-untyped]
    except ImportError:
        return text

    doc = Document(str(src))

    heading_map: dict[str, int] = {}
    for para in doc.paragraphs:
        style_name = para.style.name if para.style else ""
        if not style_name.startswith("Heading"):
            continue
        parts = style_name.split()
        if len(parts) >= 2 and parts[-1].isdigit():
            level = int(parts[-1])
            para_text = para.text.strip()
            if para_text:
                heading_map[para_text] = level

    if not heading_map:
        return text

    lines = text.split("\n")
    result = []
    for line in lines:
        stripped = line.strip()
        if not stripped:
            result.append(line)
            continue

        existing_match = re.match(r'^(#{1,6})\s+(.+)$', stripped)
        if existing_match:
            heading_text = existing_match.group(2).strip()
            if heading_text in heading_map:
                level = heading_map[heading_text]
                result.append(f"{'#' * level} {heading_text}")
                continue

        if stripped in heading_map:
            level = heading_map[stripped]
            result.append(f"{'#' * level} {stripped}")
            continue

        result.append(line)
    return "\n".join(result)


def _restore_pptx_hyperlinks(src: Path, text: str) -> str:
    """Restore hyperlinks lost by MarkItDown in pptx conversion.

    python-pptx reads each slide's text frames and finds runs that have
    a hyperlink address.  The plain text is then replaced with
    ``[text](url)`` in the Markdown output.
    """
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ImportError:
        return text

    prs = Presentation(str(src))

    links: list[tuple[str, str]] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.hyperlink and run.hyperlink.address:
                        display = run.text.strip()
                        url = run.hyperlink.address
                        if display and url:
                            links.append((display, url))

    if not links:
        return text

    for display, url in sorted(links, key=lambda x: len(x[0]), reverse=True):
        md_link = f"[{display}]({url})"
        if md_link in text:
            continue
        if f"[{display}](" in text:
            continue
        text = text.replace(display, md_link)

    return text


def _add_pptx_slide_separators(text: str) -> str:
    """Insert horizontal rules (---) between pptx slides.

    MarkItDown outputs ``<!-- Slide number: N -->`` comments.
    This function adds ``---`` before each slide comment (except the first)
    to visually separate slides in the Markdown output.
    """
    lines = text.split("\n")
    result: list[str] = []
    first_slide = True
    for line in lines:
        if line.strip().startswith("<!-- Slide number:"):
            if first_slide:
                first_slide = False
            else:
                if result and result[-1].strip():
                    result.append("")
                result.append("---")
                result.append("")
        result.append(line)
    return "\n".join(result)


def _correct_xlsx_merged_cells(src: Path, text: str) -> str:
    """Fix NaN values in xlsx tables caused by merged cells.

    MarkItDown renders merged cells as ``NaN``.  This function reads the
    workbook with openpyxl to find merged cell ranges and replaces the
    ``NaN`` placeholders with the actual merged cell value.
    It also ensures each sheet has an ``## SheetName`` heading.
    """
    try:
        from openpyxl import load_workbook  # type: ignore[import-untyped]
    except ImportError:
        return text

    wb = load_workbook(str(src), data_only=True)

    sheet_tables: dict[str, list[list[str]]] = {}
    for ws in wb.worksheets:
        merged_values: dict[tuple[int, int], str] = {}
        for merge_range in ws.merged_cells.ranges:
            top_left = ws.cell(merge_range.min_row, merge_range.min_col)
            value = str(top_left.value) if top_left.value is not None else ""
            for row in range(merge_range.min_row, merge_range.max_row + 1):
                for col in range(merge_range.min_col, merge_range.max_col + 1):
                    if (row, col) != (merge_range.min_row, merge_range.min_col):
                        merged_values[(row, col)] = value

        rows: list[list[str]] = []
        for row_idx, row in enumerate(ws.iter_rows(values_only=False), start=1):
            cells: list[str] = []
            for col_idx, cell in enumerate(row, start=1):
                if (row_idx, col_idx) in merged_values:
                    cells.append(merged_values[(row_idx, col_idx)])
                else:
                    cells.append(str(cell.value) if cell.value is not None else "")
            rows.append(cells)
        sheet_tables[ws.title] = rows

    wb.close()

    if not sheet_tables:
        return text

    lines = text.split("\n")
    result: list[str] = []

    current_sheet: str | None = None
    current_sheet_rows: list[list[str]] = []
    table_row_idx = 0
    in_separator = False

    for line in lines:
        stripped = line.strip()

        if stripped.startswith("## "):
            sheet_name = stripped[3:].strip()
            if sheet_name in sheet_tables:
                current_sheet = sheet_name
                current_sheet_rows = sheet_tables[sheet_name]
                table_row_idx = 0
                in_separator = False

        if stripped.startswith("|") and stripped.endswith("|") and current_sheet:
            cells = [c.strip() for c in stripped.split("|")[1:-1]]

            if all(c.replace("-", "").replace(":", "").strip() == "" for c in cells):
                in_separator = True
                result.append(line)
                continue

            if not in_separator:
                table_row_idx = 0
            else:
                table_row_idx += 1

            has_nan = any("NaN" in c for c in cells)
            if has_nan and current_sheet_rows:
                sheet_row_idx = 0 if not in_separator else table_row_idx
                if sheet_row_idx < len(current_sheet_rows):
                    sheet_row = current_sheet_rows[sheet_row_idx]
                    new_cells: list[str] = []
                    for i, c in enumerate(cells):
                        if "NaN" in c and i < len(sheet_row):
                            new_cells.append(sheet_row[i])
                        else:
                            new_cells.append(c)
                    line = "| " + " | ".join(new_cells) + " |"

        result.append(line)

    return "\n".join(result)
