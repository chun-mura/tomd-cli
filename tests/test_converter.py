"""Tests for tomd conversion logic."""

import pytest
from tomd.pdf import (
    _clean_cell,
    _detect_table_lines,
    _table_to_markdown,
    _tables_have_same_header,
    _merge_continuation_row,
    _collect_table_cell_texts,
    _find_table_region,
    _strip_page_headers,
    _apply_headings,
    _apply_bullets,
)
from tomd.images import (
    strip_base64_images,
    _extract_images,
    _replace_image_placeholders,
)
from tomd.office import (
    _correct_docx_headings,
    _restore_pptx_hyperlinks,
    _add_pptx_slide_separators,
    _correct_xlsx_merged_cells,
)


class TestStripBase64Images:
    def test_replaces_base64_png(self):
        text = '![alt text](data:image/png;base64,iVBORw0KGgoAAAANSUhEUg==)'
        result = strip_base64_images(text)
        assert result == '![alt text]()'

    def test_replaces_base64_jpeg(self):
        text = '![photo](data:image/jpeg;base64,/9j/4AAQSkZJRg==)'
        result = strip_base64_images(text)
        assert result == '![photo]()'

    def test_preserves_empty_alt(self):
        text = '![](data:image/png;base64,abc123)'
        result = strip_base64_images(text)
        assert result == '![]()'

    def test_preserves_normal_images(self):
        text = '![alt](https://example.com/image.png)'
        result = strip_base64_images(text)
        assert result == '![alt](https://example.com/image.png)'

    def test_preserves_plain_text(self):
        text = 'Hello world'
        result = strip_base64_images(text)
        assert result == 'Hello world'

    def test_multiple_images(self):
        text = (
            '![a](data:image/png;base64,AAA) middle '
            '![b](data:image/gif;base64,BBB)'
        )
        result = strip_base64_images(text)
        assert result == '![a]() middle ![b]()'


class TestCleanCell:
    def test_none(self):
        assert _clean_cell(None) == ""

    def test_newlines_replaced(self):
        assert _clean_cell("hello\nworld") == "hello world"

    def test_null_byte_replaced(self):
        assert _clean_cell("GPT\x004o") == "GPT-4o"

    def test_pipe_escaped(self):
        assert _clean_cell("a|b") == "a\\|b"

    def test_strips_whitespace(self):
        assert _clean_cell("  text  ") == "text"


class TestDetectTableLines:
    def test_vertical_thin_rects(self):
        rects = [
            {"x0": 10.0, "x1": 10.5, "top": 100.0, "bottom": 300.0},
            {"x0": 200.0, "x1": 200.5, "top": 100.0, "bottom": 300.0},
        ]
        v, h = _detect_table_lines(rects)
        assert len(v) == 2
        assert 100.0 in h and 300.0 in h

    def test_horizontal_thin_rects(self):
        rects = [
            {"x0": 10.0, "x1": 200.0, "top": 100.0, "bottom": 100.5},
        ]
        v, h = _detect_table_lines(rects)
        assert len(v) == 0
        assert len(h) == 1

    def test_cell_rects(self):
        rects = [
            {"x0": 10.0, "x1": 100.0, "top": 50.0, "bottom": 80.0},
        ]
        v, h = _detect_table_lines(rects)
        assert 50.0 in h and 80.0 in h


class TestTableToMarkdown:
    def test_basic_table(self):
        table = [
            ["A", "B"],
            ["1", "2"],
            ["3", "4"],
        ]
        result = _table_to_markdown(table)
        assert "| A | B |" in result
        assert "| 1 | 2 |" in result
        assert "| 3 | 4 |" in result

    def test_skips_empty_rows(self):
        table = [
            ["A", "B"],
            [None, None],
            ["1", "2"],
        ]
        result = _table_to_markdown(table)
        lines = result.strip().split("\n")
        assert len(lines) == 3  # header + separator + 1 data row

    def test_pads_short_rows(self):
        table = [
            ["A", "B", "C"],
            ["1"],
        ]
        result = _table_to_markdown(table)
        assert "| 1 |  |  |" in result

    def test_empty_table(self):
        assert _table_to_markdown([]) == ""
        assert _table_to_markdown([["A"]]) == ""


class TestTablesHaveSameHeader:
    def test_same_header(self):
        a = [["A", "B"], ["1", "2"]]
        b = [["A", "B"], ["3", "4"]]
        assert _tables_have_same_header(a, b)

    def test_different_header(self):
        a = [["A", "B"], ["1", "2"]]
        b = [["X", "Y"], ["3", "4"]]
        assert not _tables_have_same_header(a, b)

    def test_single_column_not_merged(self):
        a = [["A"], ["1"]]
        b = [["A"], ["2"]]
        assert not _tables_have_same_header(a, b)

    def test_empty_tables(self):
        assert not _tables_have_same_header([], [["A"]])


class TestMergeContinuationRow:
    def test_merges_non_empty_cells(self):
        last = ["hello", "world"]
        cont = [None, " again"]
        result = _merge_continuation_row(last, cont)
        assert result == ["hello", "world  again"]

    def test_fills_empty_cells(self):
        last = [None, "existing"]
        cont = ["new", None]
        result = _merge_continuation_row(last, cont)
        assert result == ["new", "existing"]


class TestCollectTableCellTexts:
    def test_collects_cell_lines(self):
        tables = [[["Header A", "Header B"], ["cell\nwith newline", "short"]]]
        result = _collect_table_cell_texts(tables)
        assert "Header A" in result
        assert "cell" in result
        assert "with newline" in result
        assert "short" in result

    def test_skips_empty_cells(self):
        tables = [[[None, ""], ["text", None]]]
        result = _collect_table_cell_texts(tables)
        assert "text" in result
        assert "" not in result

    def test_replaces_null_bytes(self):
        tables = [[["GPT\x004o", "other"]]]
        result = _collect_table_cell_texts(tables)
        assert "GPT-4o" in result


class TestFindTableRegion:
    def test_finds_region(self):
        lines = [
            "Title",
            "",
            "Header A",
            "cell data",
            "more cell",
            "",
            "After table text",
        ]
        cell_texts = {"Header A", "cell data", "more cell"}
        start, end = _find_table_region(lines, cell_texts)
        assert start == 2
        assert end == 4

    def test_no_match(self):
        lines = ["Title", "No table here"]
        cell_texts = {"something else entirely"}
        start, end = _find_table_region(lines, cell_texts)
        assert start is None
        assert end is None


class TestStripPageHeaders:
    def test_removes_repeated_lines(self):
        text = "Header\nContent A\nHeader\nContent B\nHeader"
        result = _strip_page_headers(text)
        assert "Header" not in result
        assert "Content A" in result

    def test_removes_page_numbers(self):
        # Header text must appear 3+ times to be detected as a page header
        text = "MyHeader\nLine 1\nMyHeader\nLine 2\nMyHeader\nLine 3\nLine 4\nMyHeader2"
        result = _strip_page_headers(text)
        assert "MyHeader2" not in result

    def test_cleans_header_prefix(self):
        text = "MyHeader\nLine 1\nMyHeader\nLine 2\nMyHeader\nLine 3\nLine 4\nMyHeader2max_tokens（最大出力量）"
        result = _strip_page_headers(text)
        assert "max_tokens（最大出力量）" in result

    def test_preserves_short_text(self):
        text = "Short"
        result = _strip_page_headers(text)
        assert result == "Short"


class TestApplyHeadings:
    def test_applies_h1_and_h2(self):
        heading_map = {"Title": 1, "Section": 2}
        text = "Title\n\nSection\n\nBody text"
        result = _apply_headings(text, heading_map)
        assert "# Title" in result
        assert "## Section" in result
        assert "Body text" in result

    def test_no_headings(self):
        text = "Just body text"
        result = _apply_headings(text, {})
        assert result == text

    def test_does_not_affect_non_matching_lines(self):
        heading_map = {"Title": 1}
        text = "Title\nNot a heading\nMore text"
        result = _apply_headings(text, heading_map)
        assert result == "# Title\nNot a heading\nMore text"


class TestApplyBullets:
    def test_exact_match(self):
        items = {"temperature", "top-k"}
        text = "temperature\n\nbody text\n\ntop-k"
        result = _apply_bullets(text, items)
        assert "- temperature" in result
        assert "- top-k" in result
        assert "body text" in result
        assert "- body" not in result

    def test_startswith_match(self):
        items = {"temperature"}
        text = "temperature（温度）\n\nbody text"
        result = _apply_bullets(text, items)
        assert "- temperature（温度）" in result

    def test_skips_parameter_values(self):
        items = {"temperature"}
        text = "temperature=0.2, top_p=0.9"
        result = _apply_bullets(text, items)
        assert not result.startswith("- ")

    def test_skips_headings_and_tables(self):
        items = {"Section"}
        text = "## Section\n| Section | col |"
        result = _apply_bullets(text, items)
        assert "## Section" in result
        assert "- ## Section" not in result
        assert "- | Section" not in result

    def test_no_items(self):
        text = "Just text"
        result = _apply_bullets(text, set())
        assert result == text


class TestReplaceImagePlaceholders:
    def test_replaces_empty_parens(self):
        image_map = {"image1.png": "images/doc_image1.png"}
        text = "Some text\n![alt]()\nMore text"
        result = _replace_image_placeholders(text, image_map)
        assert "![alt](images/doc_image1.png)" in result

    def test_replaces_multiple(self):
        image_map = {
            "image1.png": "images/doc_image1.png",
            "image2.jpg": "images/doc_image2.jpg",
        }
        text = "![first]()\n![second]()"
        result = _replace_image_placeholders(text, image_map)
        assert "![first](images/doc_image1.png)" in result
        assert "![second](images/doc_image2.jpg)" in result

    def test_no_images(self):
        text = "No images here"
        result = _replace_image_placeholders(text, {})
        assert result == text

    def test_skips_emf_wmf(self):
        image_map = {"image1.emf": "images/doc_image1.emf"}
        text = "![]()"
        result = _replace_image_placeholders(text, image_map)
        # EMF/WMF are skipped, so placeholder stays
        assert result == "![]()"

    def test_preserves_non_empty_refs(self):
        image_map = {"image1.png": "images/doc_image1.png"}
        text = "![alt](https://example.com/img.png)"
        result = _replace_image_placeholders(text, image_map)
        assert result == text

    def test_stops_when_images_exhausted(self):
        image_map = {"image1.png": "images/doc_image1.png"}
        text = "![a]()\n![b]()\n![c]()"
        result = _replace_image_placeholders(text, image_map)
        assert "![a](images/doc_image1.png)" in result
        # Remaining placeholders stay unchanged
        assert "![b]()" in result
        assert "![c]()" in result

    def test_replaces_pptx_style_refs(self):
        image_map = {"image1.png": "images/slides_image1.png"}
        text = "![slide_img.png](Picture3.jpg)"
        result = _replace_image_placeholders(text, image_map)
        assert "![slide_img.png](images/slides_image1.png)" in result

    def test_preserves_urls_in_pptx_mode(self):
        image_map = {"image1.png": "images/doc_image1.png"}
        text = "![]()\n![alt](https://example.com/img.png)"
        result = _replace_image_placeholders(text, image_map)
        assert "![](images/doc_image1.png)" in result
        assert "![alt](https://example.com/img.png)" in result


class TestExtractImages:
    def test_non_docx_pptx_returns_empty(self, tmp_path):
        pdf = tmp_path / "test.pdf"
        pdf.write_bytes(b"fake pdf")
        dest = tmp_path / "test.md"
        assert _extract_images(pdf, dest) == {}

    def test_bad_zip_returns_empty(self, tmp_path):
        docx = tmp_path / "test.docx"
        docx.write_bytes(b"not a zip file")
        dest = tmp_path / "test.md"
        assert _extract_images(docx, dest) == {}

    def test_extracts_from_docx(self, tmp_path):
        import zipfile
        docx_path = tmp_path / "test.docx"
        with zipfile.ZipFile(str(docx_path), "w") as zf:
            zf.writestr("word/media/image1.png", b"\x89PNG fake image")
            zf.writestr("word/media/image2.jpg", b"\xff\xd8 fake jpeg")
            zf.writestr("word/document.xml", b"<doc/>")
        dest = tmp_path / "output" / "test.md"
        dest.parent.mkdir(parents=True, exist_ok=True)
        result = _extract_images(docx_path, dest)
        assert "image1.png" in result
        assert "image2.jpg" in result
        assert (dest.parent / "images" / "test_image1.png").exists()
        assert (dest.parent / "images" / "test_image2.jpg").exists()

    def test_extracts_from_pptx(self, tmp_path):
        import zipfile
        pptx_path = tmp_path / "slides.pptx"
        with zipfile.ZipFile(str(pptx_path), "w") as zf:
            zf.writestr("ppt/media/image1.png", b"\x89PNG fake image")
        dest = tmp_path / "slides.md"
        result = _extract_images(pptx_path, dest)
        assert "image1.png" in result
        assert (dest.parent / "images" / "slides_image1.png").exists()


class TestCorrectDocxHeadings:
    def test_corrects_heading_levels(self, tmp_path):
        from docx import Document
        docx_path = tmp_path / "test.docx"
        doc = Document()
        doc.add_heading("Main Title", level=1)
        doc.add_paragraph("Some body text.")
        doc.add_heading("Sub Section", level=2)
        doc.add_paragraph("More body text.")
        doc.add_heading("Deep Section", level=3)
        doc.save(str(docx_path))

        # Simulate MarkItDown output with wrong heading levels
        text = "## Main Title\n\nSome body text.\n\n### Sub Section\n\nMore body text.\n\n#### Deep Section"
        result = _correct_docx_headings(docx_path, text)
        assert "# Main Title" in result
        assert "## Sub Section" in result
        assert "### Deep Section" in result

    def test_adds_missing_headings(self, tmp_path):
        from docx import Document
        docx_path = tmp_path / "test.docx"
        doc = Document()
        doc.add_heading("Title", level=1)
        doc.add_paragraph("Body text.")
        doc.save(str(docx_path))

        # MarkItDown missed the heading entirely
        text = "Title\n\nBody text."
        result = _correct_docx_headings(docx_path, text)
        assert "# Title" in result
        assert "Body text." in result

    def test_no_headings_in_docx(self, tmp_path):
        from docx import Document
        docx_path = tmp_path / "test.docx"
        doc = Document()
        doc.add_paragraph("Just a paragraph.")
        doc.save(str(docx_path))

        text = "Just a paragraph."
        result = _correct_docx_headings(docx_path, text)
        assert result == text


class TestRestorePptxHyperlinks:
    def test_restores_link(self, tmp_path):
        from pptx import Presentation
        pptx_path = tmp_path / "links.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Test"
        tf = slide.placeholders[1].text_frame
        p = tf.paragraphs[0]
        p.clear()
        run = p.add_run()
        run.text = "Visit Example"
        run.hyperlink.address = "https://example.com"
        prs.save(str(pptx_path))

        text = "# Test\nVisit Example"
        result = _restore_pptx_hyperlinks(pptx_path, text)
        assert "[Visit Example](https://example.com)" in result

    def test_preserves_existing_links(self, tmp_path):
        from pptx import Presentation
        pptx_path = tmp_path / "links.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Test"
        tf = slide.placeholders[1].text_frame
        p = tf.paragraphs[0]
        p.clear()
        run = p.add_run()
        run.text = "Click here"
        run.hyperlink.address = "https://example.com"
        prs.save(str(pptx_path))

        # Text already has the link
        text = "# Test\n[Click here](https://example.com)"
        result = _restore_pptx_hyperlinks(pptx_path, text)
        assert result.count("[Click here](https://example.com)") == 1

    def test_no_hyperlinks(self, tmp_path):
        from pptx import Presentation
        pptx_path = tmp_path / "plain.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Title"
        slide.placeholders[1].text = "No links here"
        prs.save(str(pptx_path))

        text = "# Title\nNo links here"
        result = _restore_pptx_hyperlinks(pptx_path, text)
        assert result == text


class TestAddPptxSlideSeparators:
    def test_adds_separators(self):
        text = (
            "<!-- Slide number: 1 -->\n"
            "# Slide One\n"
            "Content 1\n"
            "<!-- Slide number: 2 -->\n"
            "# Slide Two\n"
            "Content 2\n"
            "<!-- Slide number: 3 -->\n"
            "# Slide Three\n"
            "Content 3"
        )
        result = _add_pptx_slide_separators(text)
        assert result.count("---") == 2
        # First slide should NOT have a separator before it
        lines = result.split("\n")
        first_slide_idx = next(
            i for i, l in enumerate(lines)
            if "Slide number: 1" in l
        )
        # No --- before the first slide
        assert "---" not in "\n".join(lines[:first_slide_idx])

    def test_single_slide(self):
        text = "<!-- Slide number: 1 -->\n# Only Slide\nContent"
        result = _add_pptx_slide_separators(text)
        assert "---" not in result

    def test_no_slide_comments(self):
        text = "Just regular text\nNo slides here"
        result = _add_pptx_slide_separators(text)
        assert result == text


class TestCorrectXlsxMergedCells:
    def test_fixes_nan_from_merged_cells(self, tmp_path):
        from openpyxl import Workbook
        xlsx_path = tmp_path / "merged.xlsx"
        wb = Workbook()
        ws = wb.active
        ws.title = "Data"
        ws.append(["Category", "Value"])
        ws.append(["Electronics", 100])
        ws.append(["", 200])
        ws.merge_cells("A2:A3")
        ws.append(["Clothing", 50])
        wb.save(str(xlsx_path))

        # Simulate MarkItDown output with NaN
        text = (
            "## Data\n"
            "| Category | Value |\n"
            "| --- | --- |\n"
            "| Electronics | 100 |\n"
            "| NaN | 200 |\n"
            "| Clothing | 50 |"
        )
        result = _correct_xlsx_merged_cells(xlsx_path, text)
        assert "NaN" not in result
        assert "Electronics" in result
        # The merged cell should show the parent value
        lines = result.split("\n")
        row4 = [l for l in lines if "200" in l][0]
        assert "Electronics" in row4

    def test_no_merged_cells(self, tmp_path):
        from openpyxl import Workbook
        xlsx_path = tmp_path / "normal.xlsx"
        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws.append(["A", "B"])
        ws.append([1, 2])
        wb.save(str(xlsx_path))

        text = (
            "## Sheet1\n"
            "| A | B |\n"
            "| --- | --- |\n"
            "| 1 | 2 |"
        )
        result = _correct_xlsx_merged_cells(xlsx_path, text)
        assert result == text

    def test_multiple_sheets(self, tmp_path):
        from openpyxl import Workbook
        xlsx_path = tmp_path / "multi.xlsx"
        wb = Workbook()
        ws1 = wb.active
        ws1.title = "Sheet A"
        ws1.append(["X", "Y"])
        ws1.append(["a", "b"])
        ws1.append(["", "c"])
        ws1.merge_cells("A2:A3")

        ws2 = wb.create_sheet("Sheet B")
        ws2.append(["P", "Q"])
        ws2.append(["p1", "q1"])
        wb.save(str(xlsx_path))

        text = (
            "## Sheet A\n"
            "| X | Y |\n"
            "| --- | --- |\n"
            "| a | b |\n"
            "| NaN | c |\n"
            "\n"
            "## Sheet B\n"
            "| P | Q |\n"
            "| --- | --- |\n"
            "| p1 | q1 |"
        )
        result = _correct_xlsx_merged_cells(xlsx_path, text)
        assert "NaN" not in result
        # Sheet A merged cell should be filled
        lines = result.split("\n")
        nan_line = [l for l in lines if "| c |" in l][0]
        assert "| a | c |" in nan_line
        # Sheet B should be unchanged
        assert "| p1 | q1 |" in result
